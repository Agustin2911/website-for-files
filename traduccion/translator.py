from googletrans import Translator
from docx import Document
import os
import shutil
from pptx import * 
from docx2pdf import convert
from pdf2docx import Converter

#translator
def translate_text(text, target_language='en'):
    try:
        translation = translator.translate(text, dest=target_language)
        print(translation.text)
        return translation.text
    except Exception as e:
        print(f'error :${e}')
        return " "

#convetors
    
#convetor of pdf to docx
def pdf_word(pdf_path, word_path):
    cv = Converter(pdf_path)
    cv.convert(word_path, start=0, end=None)
    cv.close()
    os.remove(pdf_path)

#convertor of docx to pdf
def docx_pdf(docx_path, pdf_path):
    try:
        convert(docx_path, pdf_path)
    except Exception as e:
        print(f"Error durante la conversión: {str(e)}")
    os.remove(docx_path)



#readers of files

#reader of docx
def read_docx(input_path, output_path,language):
    doc = Document(input_path)
    for paragraph in doc.paragraphs:
        if paragraph.text:
            paragraph.text = translate_text(paragraph.text,language)

    doc.save(output_path)

#reader for the case of pdf
def read_docx2(input_path, output_path,language):
    doc = Document(input_path)
    for paragraph in doc.paragraphs:
        if paragraph.text:
            paragraph.text = translate_text(paragraph.text,language)

    doc.save(output_path)
    os.remove(input_path)

#reader of pptx
def read_pptx(input_path, output_path, language):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                translated_text = translate_text(shape.text,language)
                if translated_text!=" ":
                    formato_texto=shape.text_frame.paragraphs[0].runs[0].font
                    shape.text = translated_text
                    shape.text_frame.paragraphs[0].runs[0].font.name=formato_texto.name
                    shape.text_frame.paragraphs[0].runs[0].font.size= formato_texto.size
     
    prs.save(output_path)

#reader of ppt
def read_ppt(input_path, output_path, target_language='en'):
    prs = Presentation(input_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                translated_text = translate_text(shape.text, target_language)
                if translated_text==" ":
                    formato_texto=shape.text_frame.paragraphs[0].runs[0].font
                    shape.text = translated_text
                    shape.text_frame.paragraphs[0].runs[0].font.name=formato_texto.name
                    shape.text_frame.paragraphs[0].runs[0].font.size= formato_texto.size

    prs.save(output_path)

#copyer of arch
def copy_arch(origen, destino):
    try:
        shutil.copy(origen, destino)
        print(f"Archivo copiado de {origen} a {destino} exitosamente.")
    except FileNotFoundError:
        print(f"Error: El archivo {origen} no fue encontrado.")
    except PermissionError:
        print(f"Error: No tienes permisos para copiar el archivo.")


#main program

translator=Translator()

